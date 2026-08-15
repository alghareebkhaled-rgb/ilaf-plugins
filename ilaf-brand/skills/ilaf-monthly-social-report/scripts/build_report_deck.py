#!/usr/bin/env python3
"""
Build the ILAF monthly social performance deck from a data JSON file.

    python3 build_report_deck.py report_data.json -o ILAF_Social_Report_2026-07.pptx

The whole point of this script is that the monthly report should cost minutes, not hours.
It owns the brand styling, the bilingual layout, the RTL handling and the charts, so the
analysis work upstream can stay focused on the numbers and the story they tell.

Missing values are rendered as "غير متاح · Not available" rather than zero. That distinction
carries real meaning in this report: zero means they published nothing, unavailable means we
could not see. Collapsing the two misleads the reader.

If you need a slide or a chart this script doesn't do, extend it here rather than hand-building
slides for one month — a script that drifts out of date stops saving anybody time.
"""

import argparse
import json
import os
import sys
from copy import deepcopy

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------------------
# Brand
# --------------------------------------------------------------------------------------

MAGENTA = RGBColor(0xA5, 0x0B, 0x7F)
NAVY = RGBColor(0x00, 0x1E, 0x6A)
GOLD = RGBColor(0xFD, 0xB7, 0x2F)
CREAM = RGBColor(0xF0, 0xCF, 0x94)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x6B, 0x6B, 0x80)
LIGHT = RGBColor(0xF5, 0xF3, 0xF7)
GOOD = RGBColor(0x1B, 0x8A, 0x5A)
BAD = RGBColor(0xC0, 0x39, 0x2B)

# Google Slides has to be able to render these, so we use fonts that exist there.
# The bundled brand fonts (Bahij, JF Flat) can't be embedded in a pptx and would silently
# fall back to something arbitrary — Tajawal is the brand's own documented Arabic fallback.
FONT_AR = "Tajawal"
FONT_EN = "Arial"

NA = "غير متاح · Not available"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(os.path.dirname(HERE), "assets")
LOGO = os.path.join(ASSETS, "logo", "ilaf-logo-white.png")
SADU = os.path.join(ASSETS, "elements", "sadu-band-gold.png")

CHART_DIR = None  # set at runtime


# --------------------------------------------------------------------------------------
# Small helpers
# --------------------------------------------------------------------------------------

def has(v):
    """A value is present if it isn't None and isn't an empty string."""
    return v is not None and v != ""


def fmt_num(v):
    if not has(v):
        return NA
    if isinstance(v, float):
        if abs(v - round(v)) < 1e-9:
            return f"{int(round(v)):,}"
        return f"{v:,.1f}"
    if isinstance(v, int):
        return f"{v:,}"
    return str(v)


def fmt_pct(v, decimals=1):
    if not has(v):
        return NA
    return f"{float(v):.{decimals}f}%"


def delta(cur, prev):
    """Return (text, pct_change) or (None, None) when it can't be computed."""
    if not has(cur) or not has(prev):
        return None, None
    try:
        cur, prev = float(cur), float(prev)
    except (TypeError, ValueError):
        return None, None
    if prev == 0:
        return None, None
    pct = (cur - prev) / abs(prev) * 100.0
    arrow = "▲" if pct > 0 else ("▼" if pct < 0 else "▬")
    return f"{arrow} {abs(pct):.1f}%", pct


def delta_color(pct, good_direction="up"):
    """Colour by whether the movement is GOOD, not by whether it is UP.

    Unfollows rising is an up arrow and a bad month. Getting this backwards on the
    summary slide is the most damaging error the deck can make.
    """
    if pct is None or good_direction == "neutral":
        return MUTED
    if abs(pct) < 0.05:
        return MUTED
    improving = pct > 0 if good_direction == "up" else pct < 0
    return GOOD if improving else BAD


def set_rtl(paragraph):
    """python-pptx has no RTL property, so set it on the underlying XML."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    return tb, tf


def write(tf, text, size=14, color=INK, bold=False, arabic=False,
          align=None, first=False, space_after=0, italic=False):
    """Append a paragraph (or fill the first one) with brand-consistent styling."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text or ""
    p.space_after = Pt(space_after)
    if align is not None:
        p.alignment = align
    elif arabic:
        p.alignment = PP_ALIGN.RIGHT
    if arabic:
        set_rtl(p)
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = FONT_AR if arabic else FONT_EN
        if arabic:
            set_cs_font(run, FONT_AR)
    return p


def set_cs_font(run, typeface):
    """Set the complex-script typeface on an Arabic run.

    python-pptx's font.name only writes <a:latin>, which governs Latin characters.
    Arabic is complex script and reads <a:cs> — without it the renderer silently
    substitutes whatever it likes, which is how Arabic text ends up as empty boxes in
    some viewers while looking fine in others. Setting both keeps the two in step
    wherever the file travels.
    """
    rPr = run._r.get_or_add_rPr()
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for tag in ("cs", "ea"):
        el = rPr.find(f"{ns}{tag}")
        if el is None:
            el = rPr.makeelement(f"{ns}{tag}", {})
            rPr.append(el)
        el.set("typeface", typeface)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def rounded(slide, x, y, w, h, fill=None, line=None, line_w=1.0, adj=0.08):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = adj
    except (IndexError, AttributeError):
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


# --------------------------------------------------------------------------------------
# Slide chrome
# --------------------------------------------------------------------------------------

DATA_STATUS = "verified"   # set from the data file at runtime

SAMPLE_BANNER_EN = "SAMPLE DATA — NOT REAL FIGURES"
SAMPLE_BANNER_AR = "بيانات تجريبية — أرقام غير حقيقية"


def base_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return slide


def sample_stamp(slide, on_dark=False):
    """Mark every slide when the figures aren't from a verified source.

    A deck built from placeholder or example numbers looks exactly as finished as a real
    one, and decks get forwarded. The stamp has to be on the slide itself — a caveat in
    the covering message doesn't travel with the file.
    """
    if DATA_STATUS == "verified":
        return
    rect(slide, Emu(0), Inches(7.16), SLIDE_W, Inches(0.34), fill=BAD)
    _, tf = textbox(slide, Inches(0.5), Inches(7.21), Inches(6.0), Inches(0.26))
    write(tf, SAMPLE_BANNER_EN, size=9.5, color=WHITE, bold=True, first=True)
    _, tf2 = textbox(slide, Inches(7.0), Inches(7.21), Inches(5.8), Inches(0.26))
    write(tf2, SAMPLE_BANNER_AR, size=9.5, color=WHITE, bold=True, arabic=True, first=True)


PAGE = {"n": 1, "total": None}


def header(slide, title_en, title_ar, num=None, total=None):
    """Magenta band, gold rule under it, bilingual title.

    Page numbers count themselves. They used to be passed in by hand, which meant
    adding or paginating a slide silently mislabelled every slide after it.
    """
    rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(1.05), fill=MAGENTA)
    rect(slide, Emu(0), Inches(1.05), SLIDE_W, Pt(3), fill=GOLD)

    _, tf = textbox(slide, Inches(0.6), Inches(0.18), Inches(7.4), Inches(0.72),
                    anchor=MSO_ANCHOR.MIDDLE)
    write(tf, title_en, size=25, color=WHITE, bold=True, first=True)

    _, tf2 = textbox(slide, Inches(8.3), Inches(0.18), Inches(4.4), Inches(0.72),
                     anchor=MSO_ANCHOR.MIDDLE)
    write(tf2, title_ar, size=22, color=CREAM, bold=True, arabic=True, first=True)

    PAGE["n"] += 1
    tot = PAGE["total"]
    label = f"{PAGE['n']}/{tot}" if tot else str(PAGE["n"])
    _, tfn = textbox(slide, Inches(12.3), Inches(6.95), Inches(0.7), Inches(0.35))
    write(tfn, label, size=10, color=MUTED, align=PP_ALIGN.RIGHT, first=True)
    sample_stamp(slide)


def subtitle(slide, text_en, text_ar=None, y=1.22):
    """A plain-language line under the header, for slides whose title is a term of art.

    English and Arabic get separate text boxes rather than two paragraphs in one. Mixing
    an LTR and an RTL paragraph in a single frame renders unreliably once the file is
    converted between PowerPoint, Slides and PDF — separate frames survive every hop.
    """
    if text_en:
        _, tf = textbox(slide, Inches(0.6), Inches(y), Inches(11.9), Inches(0.32))
        write(tf, text_en, size=11.5, color=NAVY, first=True)
    if text_ar:
        _, tf2 = textbox(slide, Inches(0.6), Inches(y + 0.30), Inches(11.9), Inches(0.32))
        write(tf2, text_ar, size=11, color=MUTED, arabic=True, first=True)


def footer_note(slide, text_en, text_ar=None, y=6.85):
    _, tf = textbox(slide, Inches(0.6), Inches(y), Inches(11.4), Inches(0.45))
    write(tf, text_en, size=9.5, color=MUTED, first=True, italic=True)
    if text_ar:
        write(tf, text_ar, size=9.5, color=MUTED, arabic=True, italic=True)


# --------------------------------------------------------------------------------------
# Charts
# --------------------------------------------------------------------------------------

def _style_axes(ax):
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#CCCCCC")
    ax.spines["bottom"].set_color("#CCCCCC")
    ax.tick_params(colors="#6B6B80", labelsize=9)
    ax.grid(axis="y", color="#EEEEEE", linewidth=0.8)
    ax.set_axisbelow(True)


def save_chart(fig, name):
    """Render a chart and shrink it.

    Charts here are flat brand colours on white, so a 64-colour palette is visually
    identical and a fraction of the bytes. Deck size matters because the file gets
    uploaded to Drive and emailed around; a 500KB deck is friction nobody needs.
    """
    path = os.path.join(CHART_DIR, name)
    fig.savefig(path, dpi=110, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    try:
        from PIL import Image
        im = Image.open(path).convert("RGB")
        im.quantize(colors=64, method=Image.MEDIANCUT).save(path, optimize=True)
    except Exception:
        pass  # a bigger chart beats no chart
    return path


def chart_two_month_bars(labels, prev_vals, cur_vals, prev_label, cur_label, name):
    """Charts are English-labelled on purpose: matplotlib does not shape Arabic text
    correctly without extra dependencies, and broken Arabic glyphs look worse than English."""
    idx = range(len(labels))
    fig, ax = plt.subplots(figsize=(6.2, 3.1))
    w = 0.36
    ax.bar([i - w / 2 for i in idx], prev_vals, w, label=prev_label, color="#D9C7D6")
    ax.bar([i + w / 2 for i in idx], cur_vals, w, label=cur_label, color="#A50B7F")
    ax.set_xticks(list(idx))
    ax.set_xticklabels(labels)
    ax.legend(frameon=False, fontsize=9, loc="upper left")
    _style_axes(ax)
    return save_chart(fig, name)


def chart_hbar(labels, values, name, highlight_index=0, xlabel=""):
    fig, ax = plt.subplots(figsize=(6.0, max(2.6, 0.42 * len(labels))))
    colors = ["#A50B7F" if i == highlight_index else "#C9C4D4" for i in range(len(labels))]
    ax.barh(labels, values, color=colors, height=0.62)
    ax.invert_yaxis()
    for i, v in enumerate(values):
        ax.text(v, i, f"  {v:,.0f}", va="center", fontsize=9, color="#4A4A5A")
    ax.set_xlabel(xlabel, fontsize=9, color="#6B6B80")
    _style_axes(ax)
    ax.grid(axis="y", visible=False)
    ax.grid(axis="x", color="#EEEEEE", linewidth=0.8)
    return save_chart(fig, name)


def chart_stacked_themes(entities, themes, matrix, name):
    palette = ["#A50B7F", "#001E6A", "#FDB72F", "#F0CF94", "#7B3F9E",
               "#3D8BC4", "#E0719B", "#5CB88E", "#B0AEC0", "#8A6D3B"]
    fig, ax = plt.subplots(figsize=(7.4, 3.3))
    bottoms = [0] * len(entities)
    for ti, theme in enumerate(themes):
        vals = [matrix[ei][ti] for ei in range(len(entities))]
        ax.bar(entities, vals, bottom=bottoms, label=theme.replace("_", " ").title(),
               color=palette[ti % len(palette)], width=0.6)
        bottoms = [b + v for b, v in zip(bottoms, vals)]
    # The legend sits clear of the rotated handle labels; anything above ~-0.28 overlaps them.
    ax.legend(frameon=False, fontsize=7.5, ncol=2, loc="upper center",
              bbox_to_anchor=(0.5, -0.30))
    plt.xticks(rotation=30, ha="right", fontsize=8)
    _style_axes(ax)
    return save_chart(fig, name)


def chart_donut(share_pct, name, center_label=""):
    fig, ax = plt.subplots(figsize=(3.5, 3.5))
    rest = max(0.0, 100.0 - float(share_pct))
    ax.pie([float(share_pct), rest], colors=["#A50B7F", "#E6E2EA"],
           startangle=90, counterclock=False,
           wedgeprops=dict(width=0.34, edgecolor="white", linewidth=2))
    ax.text(0, 0.08, f"{float(share_pct):.1f}%", ha="center", va="center",
            fontsize=25, color="#A50B7F", fontweight="bold")
    ax.text(0, -0.24, center_label, ha="center", va="center", fontsize=8.5, color="#6B6B80")
    ax.set_aspect("equal")
    return save_chart(fig, name)


# --------------------------------------------------------------------------------------
# Slides
# --------------------------------------------------------------------------------------

def slide_cover(prs, d):
    slide = base_slide(prs)
    rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=MAGENTA)
    if DATA_STATUS != "verified":
        rect(slide, Inches(0.8), Inches(5.95), Inches(7.2), Inches(0.78), fill=BAD)
        _, tfw = textbox(slide, Inches(0.8), Inches(6.06), Inches(7.2), Inches(0.6))
        write(tfw, SAMPLE_BANNER_EN, size=15, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, first=True, space_after=1)
        write(tfw, SAMPLE_BANNER_AR, size=12, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, arabic=True)
    rect(slide, Emu(0), Inches(6.9), SLIDE_W, Inches(0.6), fill=NAVY)
    if os.path.exists(SADU):
        try:
            slide.shapes.add_picture(SADU, Emu(0), Inches(6.62), width=SLIDE_W)
        except Exception:
            pass
    if os.path.exists(LOGO):
        try:
            slide.shapes.add_picture(LOGO, Inches(0.75), Inches(0.6), height=Inches(1.15))
        except Exception:
            pass

    p = (d.get("period") or {})
    _, tf = textbox(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.6))
    write(tf, "Monthly Social Media Performance Report", size=40, color=WHITE,
          bold=True, first=True, space_after=4)
    write(tf, "تقرير الأداء الشهري على منصات التواصل الاجتماعي", size=28, color=CREAM,
          arabic=True, align=PP_ALIGN.LEFT, space_after=16)
    rect(slide, Inches(0.8), Inches(4.55), Inches(2.6), Pt(3), fill=GOLD)

    _, tf2 = textbox(slide, Inches(0.8), Inches(4.85), Inches(11.7), Inches(1.2))
    label = f"{p.get('month_label_en', '')}  ·  {p.get('month_label_ar', '')}"
    write(tf2, label, size=22, color=WHITE, bold=True, first=True, space_after=6)
    rng = ""
    if p.get("start") and p.get("end"):
        rng = f"{p['start']} → {p['end']}"
    gen = d.get("generated_on", "")
    sub = "   |   ".join([x for x in [rng, f"Generated {gen}" if gen else ""] if x])
    write(tf2, sub, size=12, color=CREAM)

    _, tf3 = textbox(slide, Inches(0.8), Inches(6.98), Inches(11.7), Inches(0.4))
    write(tf3, "ILAF Takaful Insurance Company · إيلاف للتأمين التكافلي",
          size=11, color=WHITE, first=True)
    return slide


def kpi_tile(slide, x, y, w, h, label_en, label_ar, value, delta_text, dcolor, sub=None):
    rounded(slide, x, y, w, h, fill=LIGHT, line=CREAM, line_w=1.0, adj=0.10)
    rect(slide, x, y, Pt(4), h, fill=MAGENTA)

    _, tf = textbox(slide, x + Inches(0.22), y + Inches(0.16), w - Inches(0.4), Inches(0.34))
    write(tf, label_en, size=10.5, color=MUTED, bold=True, first=True)

    _, tfv = textbox(slide, x + Inches(0.22), y + Inches(0.52), w - Inches(0.4), Inches(0.6))
    is_na = value == NA
    write(tfv, value, size=15 if is_na else 27, color=MUTED if is_na else NAVY,
          bold=True, first=True)

    _, tfa = textbox(slide, x + Inches(0.22), y + Inches(1.12), w - Inches(0.4), Inches(0.3))
    write(tfa, label_ar, size=11, color=MUTED, arabic=True, first=True)

    _, tfd = textbox(slide, x + Inches(0.22), y + Inches(1.44), w - Inches(0.4), Inches(0.3))
    write(tfd, delta_text or "—", size=11.5, color=dcolor, bold=True, first=True)
    if sub:
        _, tfs = textbox(slide, x + Inches(0.22), y + Inches(1.72), w - Inches(0.4), Inches(0.28))
        write(tfs, sub, size=8.5, color=MUTED, first=True)


def slide_glance(prs, d):
    slide = base_slide(prs)
    header(slide, "Month at a Glance", "نظرة عامة على الشهر", 2)

    ig = (d.get("ilaf") or {}).get("instagram") or {}
    cur, prev = (ig.get("current") or {}), (ig.get("previous") or {})
    gb = (d.get("ilaf") or {}).get("google_business") or {}
    gbc, gbp = (gb.get("current") or {}), (gb.get("previous") or {})

    tiles = [
        ("Followers", "المتابعون", cur.get("followers"), prev.get("followers"), "up", fmt_num),
        ("Reach", "الوصول", cur.get("reach"), prev.get("reach"), "up", fmt_num),
        ("Engagement Rate", "معدل التفاعل", cur.get("engagement_rate"),
         prev.get("engagement_rate"), "up", lambda v: fmt_pct(v)),
        ("Google Business Calls", "المكالمات عبر جوجل", gbc.get("phone_calls"),
         gbp.get("phone_calls"), "up", fmt_num),
        ("Posts Published", "المنشورات", cur.get("posts_total"), prev.get("posts_total"),
         "neutral", fmt_num),
    ]

    x0, y0 = Inches(0.6), Inches(1.55)
    tw, th, gap = Inches(2.33), Inches(2.05), Inches(0.19)
    for i, (le, la, c, p, gd, f) in enumerate(tiles):
        dtxt, pct = delta(c, p)
        kpi_tile(slide, x0 + i * (tw + gap), y0, tw, th, le, la,
                 f(c) if has(c) else NA, dtxt, delta_color(pct, gd),
                 sub="vs previous month" if dtxt else "no comparison available")

    hl = (d.get("headline") or {})
    rounded(slide, Inches(0.6), Inches(3.95), Inches(11.9), Inches(1.55),
            fill=WHITE, line=GOLD, line_w=1.5, adj=0.06)
    _, tf = textbox(slide, Inches(1.0), Inches(4.2), Inches(11.1), Inches(1.1))
    write(tf, hl.get("en", ""), size=14.5, color=NAVY, bold=True, first=True, space_after=8)
    write(tf, hl.get("ar", ""), size=13.5, color=INK, arabic=True)

    footer_note(slide, "Arrows show whether the movement is favourable, not simply upward.",
                "تشير الأسهم إلى اتجاه التحسن، لا إلى الارتفاع فقط.", y=5.75)
    return slide


def _metric_row(slide, y, label_en, label_ar, cur, prev, formatter, good="up"):
    rounded(slide, Inches(0.6), Inches(y), Inches(6.0), Inches(0.62),
            fill=LIGHT, line=None, adj=0.25)
    _, tf = textbox(slide, Inches(0.85), Inches(y + 0.10), Inches(2.6), Inches(0.42))
    write(tf, label_en, size=11.5, color=INK, bold=True, first=True)
    _, tfa = textbox(slide, Inches(3.35), Inches(y + 0.10), Inches(1.5), Inches(0.42))
    write(tfa, label_ar, size=10.5, color=MUTED, arabic=True, first=True)
    _, tfv = textbox(slide, Inches(4.75), Inches(y + 0.08), Inches(1.05), Inches(0.45))
    val = formatter(cur) if has(cur) else NA
    write(tfv, val, size=10 if val == NA else 13.5, color=MUTED if val == NA else NAVY,
          bold=True, align=PP_ALIGN.RIGHT, first=True)
    dtxt, pct = delta(cur, prev)
    _, tfd = textbox(slide, Inches(5.85), Inches(y + 0.13), Inches(0.65), Inches(0.4))
    write(tfd, dtxt or "—", size=10, color=delta_color(pct, good),
          bold=True, align=PP_ALIGN.RIGHT, first=True)


def slide_audience(prs, d):
    slide = base_slide(prs)
    header(slide, "Audience Growth", "نمو الجمهور", 3)
    ig = (d.get("ilaf") or {}).get("instagram") or {}
    cur, prev = (ig.get("current") or {}), (ig.get("previous") or {})
    p = (d.get("period") or {})

    rows = [
        ("Followers (end of month)", "المتابعون", "followers", fmt_num, "up"),
        ("Net growth", "صافي النمو", "follower_growth", fmt_num, "up"),
        ("New follows", "متابعات جديدة", "follows", fmt_num, "up"),
        ("Unfollows", "إلغاء المتابعة", "unfollows", fmt_num, "down"),
        ("Profile views", "زيارات الملف", "profile_views", fmt_num, "up"),
        ("Website clicks", "النقرات على الموقع", "website_clicks", fmt_num, "up"),
    ]
    y = 1.6
    for le, la, key, f, good in rows:
        _metric_row(slide, y, le, la, cur.get(key), prev.get(key), f, good)
        y += 0.72

    labels, pv, cv = [], [], []
    for le, _, key, _, _ in rows[:3]:
        if has(cur.get(key)) and has(prev.get(key)):
            labels.append(le.split(" (")[0])
            pv.append(float(prev[key]))
            cv.append(float(cur[key]))
    if labels:
        path = chart_two_month_bars(
            labels, pv, cv,
            p.get("prev_month_label_en", "Previous"), p.get("month_label_en", "Current"),
            "audience.png")
        slide.shapes.add_picture(path, Inches(7.0), Inches(1.75), width=Inches(5.7))
    else:
        _, tf = textbox(slide, Inches(7.0), Inches(3.2), Inches(5.7), Inches(0.6))
        write(tf, NA, size=13, color=MUTED, align=PP_ALIGN.CENTER, first=True)

    note = ig.get("audience_note") or {}
    if note.get("en") or note.get("ar"):
        _, tf = textbox(slide, Inches(0.6), Inches(6.1), Inches(12.0), Inches(0.75))
        if note.get("en"):
            write(tf, note["en"], size=11.5, color=NAVY, first=True)
        if note.get("ar"):
            write(tf, note["ar"], size=11, color=INK, arabic=True)
    return slide


def slide_reach(prs, d):
    slide = base_slide(prs)
    header(slide, "Reach & Engagement", "الوصول والتفاعل", 4)
    ig = (d.get("ilaf") or {}).get("instagram") or {}
    cur, prev = (ig.get("current") or {}), (ig.get("previous") or {})
    p = (d.get("period") or {})

    rows = [
        ("Reach", "الوصول", "reach", fmt_num, "up"),
        ("Impressions", "مرات الظهور", "impressions", fmt_num, "up"),
        ("Total engagement", "إجمالي التفاعل", "total_engagement", fmt_num, "up"),
        ("Likes", "الإعجابات", "likes", fmt_num, "up"),
        ("Comments", "التعليقات", "comments", fmt_num, "up"),
        ("Saves", "الحفظ", "saves", fmt_num, "up"),
        ("Shares", "المشاركات", "shares", fmt_num, "up"),
        ("Engagement rate", "معدل التفاعل", "engagement_rate", lambda v: fmt_pct(v), "up"),
    ]
    y = 1.55
    for le, la, key, f, good in rows:
        _metric_row(slide, y, le, la, cur.get(key), prev.get(key), f, good)
        y += 0.65

    labels, pv, cv = [], [], []
    for le, _, key, _, _ in [rows[0], rows[1], rows[2]]:
        if has(cur.get(key)) and has(prev.get(key)):
            labels.append(le)
            pv.append(float(prev[key]))
            cv.append(float(cur[key]))
    if labels:
        path = chart_two_month_bars(
            labels, pv, cv,
            p.get("prev_month_label_en", "Previous"), p.get("month_label_en", "Current"),
            "reach.png")
        slide.shapes.add_picture(path, Inches(7.0), Inches(1.9), width=Inches(5.7))

    note = ig.get("engagement_note") or {}
    if note.get("en") or note.get("ar"):
        _, tf = textbox(slide, Inches(7.0), Inches(5.15), Inches(5.7), Inches(1.4))
        if note.get("en"):
            write(tf, note["en"], size=11, color=NAVY, first=True, space_after=6)
        if note.get("ar"):
            write(tf, note["ar"], size=10.5, color=INK, arabic=True)

    footer_note(slide,
                "Engagement rate for ILAF is engagement ÷ reach. Competitor rates later in this "
                "deck use followers as the denominator, as competitor reach is not public.",
                None, y=6.9)
    return slide


def slide_output(prs, d):
    slide = base_slide(prs)
    header(slide, "Content Output", "حجم المحتوى المنشور", 5)
    ig = (d.get("ilaf") or {}).get("instagram") or {}
    cur, prev = (ig.get("current") or {}), (ig.get("previous") or {})

    cf = cur.get("by_format") or {}
    pf = prev.get("by_format") or {}
    formats = ["static", "carousel", "reel"]
    labels = ["Static", "Carousel", "Reel"]
    if any(has(cf.get(k)) for k in formats):
        path = chart_two_month_bars(
            labels,
            [float(pf.get(k) or 0) for k in formats],
            [float(cf.get(k) or 0) for k in formats],
            (d.get("period") or {}).get("prev_month_label_en", "Previous"),
            (d.get("period") or {}).get("month_label_en", "Current"),
            "formats.png")
        slide.shapes.add_picture(path, Inches(0.7), Inches(1.7), width=Inches(6.4))

    tiles = [
        ("Total posts", "إجمالي المنشورات", cur.get("posts_total"), prev.get("posts_total")),
        ("Stories published", "القصص المنشورة", cur.get("stories_total"), prev.get("stories_total")),
        ("Reels", "الريلز", cf.get("reel"), pf.get("reel")),
        ("Posts per week", "منشورات أسبوعياً", cur.get("posts_per_week"), prev.get("posts_per_week")),
    ]
    for i, (le, la, c, p_) in enumerate(tiles):
        col, row = i % 2, i // 2
        x = Inches(7.5) + col * Inches(2.6)
        y = Inches(1.7) + row * Inches(2.2)
        dtxt, pct = delta(c, p_)
        kpi_tile(slide, x, y, Inches(2.4), Inches(2.0), le, la,
                 fmt_num(c) if has(c) else NA, dtxt, delta_color(pct, "neutral"))

    note = ig.get("output_note") or {}
    if note.get("en") or note.get("ar"):
        # Sits below the second row of KPI tiles (which end at y=5.9) — a note placed
        # any higher runs underneath them.
        _, tf = textbox(slide, Inches(0.7), Inches(6.02), Inches(11.9), Inches(1.1))
        if note.get("en"):
            write(tf, note["en"], size=12, color=NAVY, first=True, space_after=6)
        if note.get("ar"):
            write(tf, note["ar"], size=11.5, color=INK, arabic=True)
    return slide


def _post_thumb(slide, path, x, y, box_w, box_h):
    """Place a post image inside a fixed box, cropped to fill rather than squashed.

    Instagram posts are 1:1, 4:5 and 9:16 in the same month. Scaling each to the same
    box distorts them; cropping to the box keeps faces and text upright, which is what
    makes the card recognisable at a glance.
    """
    try:
        from PIL import Image
        im = Image.open(path)
        iw, ih = im.size
    except Exception:
        return False
    box_ar = box_w / box_h
    if iw / ih > box_ar:                      # too wide — crop the sides
        new_w = int(ih * box_ar)
        im = im.crop(((iw - new_w) // 2, 0, (iw - new_w) // 2 + new_w, ih))
    else:                                     # too tall — crop top and bottom
        new_h = int(iw / box_ar)
        top = int((ih - new_h) * 0.35)        # bias upward; captions sit low
        im = im.crop((0, top, iw, top + new_h))
    # Resample to roughly what the slide actually shows. A 1080px Instagram export
    # dropped in at native size makes the deck several times bigger for detail nobody
    # can see in a 1.2-inch card.
    target_w = int(box_w * 200)               # 200dpi is past the point of visible gain
    if im.width > target_w:
        im = im.resize((target_w, max(1, int(target_w / box_ar))), Image.LANCZOS)
    out = os.path.join(CHART_DIR, "thumb_%s" % os.path.basename(path))
    im.convert("RGB").save(out, "JPEG", quality=82, optimize=True)
    slide.shapes.add_picture(out, Inches(x), Inches(y), Inches(box_w), Inches(box_h))
    rounded(slide, Inches(x), Inches(y), Inches(box_w), Inches(box_h),
            fill=None, line=CREAM, line_w=1.0, adj=0.02)
    return True


def _post_card(slide, x, y, w, h, post, accent, rank=None):
    """One post card. Shows the post image when the data file supplies `image`.

    Numbers tell the reader which post won; the picture tells them what it looked like.
    Without it, "the World Cup carousel" means nothing to anyone who wasn't in the
    room when it was published.
    """
    rounded(slide, x, y, w, h, fill=WHITE, line=accent, line_w=1.5, adj=0.05)
    rect(slide, x, y, w, Pt(5), fill=accent)

    xi, wi = x / 914400.0, w / 914400.0       # work in inches inside the card
    yi = y / 914400.0
    pad = 0.22

    img = post.get("image")
    drawn = False
    tw = 1.22                                  # thumbnail width
    th = 1.52                                  # 4:5, Instagram's feed ratio
    if img and os.path.exists(img):
        drawn = _post_thumb(slide, img, xi + pad, yi + 0.2, tw, th)

    text_x = xi + pad + (tw + 0.18 if drawn else 0)
    text_w = wi - pad * 2 - (tw + 0.18 if drawn else 0)

    head = " · ".join([str(v) for v in [
        post.get("date"), (post.get("format") or "").title(),
        (post.get("theme") or "").replace("_", " ").title()] if v])
    _, tf = textbox(slide, Inches(text_x), Inches(yi + 0.2), Inches(text_w), Inches(0.46))
    write(tf, (f"#{rank}  " if rank else "") + head, size=8.5, color=MUTED, bold=True, first=True)

    # Beside a thumbnail the title runs in a narrow column and wraps to more lines, so
    # the Arabic line and the stats row both sit lower than in the no-image layout.
    _, tf2 = textbox(slide, Inches(text_x), Inches(yi + 0.62), Inches(text_w),
                     Inches(0.62 if drawn else 0.44))
    write(tf2, post.get("topic_en", ""), size=11, color=NAVY, bold=True, first=True)
    if post.get("topic_ar"):
        _, tf2b = textbox(slide, Inches(text_x), Inches(yi + (1.28 if drawn else 1.12)),
                          Inches(text_w), Inches(0.42))
        write(tf2b, post["topic_ar"], size=10, color=INK, arabic=True, first=True)

    stats = []
    if has(post.get("reach")):
        stats.append(f"Reach {fmt_num(post['reach'])}")
    if has(post.get("engagement")):
        stats.append(f"Engagement {fmt_num(post['engagement'])}")
    if has(post.get("engagement_rate")):
        stats.append(f"ER {fmt_pct(post['engagement_rate'])}")
    sy = (yi + 0.2 + th + 0.14) if drawn else (yi + 1.66)
    _, tf3 = textbox(slide, Inches(xi + pad), Inches(sy), Inches(wi - pad * 2), Inches(0.3))
    write(tf3, "   ".join(stats) if stats else NA, size=10.5,
          color=accent if stats else MUTED, bold=True, first=True)

    _, tf4 = textbox(slide, Inches(xi + pad), Inches(sy + 0.32),
                     Inches(wi - pad * 2), Inches(1.0))
    if post.get("why_en"):
        write(tf4, post["why_en"], size=9.5, color=INK, first=True, space_after=4)
    if post.get("why_ar"):
        write(tf4, post["why_ar"], size=9, color=MUTED, arabic=True)


def slide_top_posts(prs, d):
    slide = base_slide(prs)
    header(slide, "Top Performing Posts", "أفضل المنشورات أداءً", 6)
    ig = (d.get("ilaf") or {}).get("instagram") or {}
    posts = (ig.get("top_posts") or [])[:3]
    if not posts:
        _, tf = textbox(slide, Inches(0.6), Inches(3.3), Inches(12.1), Inches(0.6))
        write(tf, NA, size=15, color=MUTED, align=PP_ALIGN.CENTER, first=True)
        footer_note(slide, "Per-post data was not available this month.",
                    "لم تتوفر بيانات المنشورات الفردية هذا الشهر.")
        return slide
    x, w = Inches(0.6), Inches(3.9)
    for i, post in enumerate(posts):
        _post_card(slide, x + i * (w + Inches(0.2)), Inches(1.5), w, Inches(3.25),
                   post, MAGENTA, rank=i + 1)
    pat = ig.get("top_pattern") or {}
    if pat.get("en") or pat.get("ar"):
        rounded(slide, Inches(0.6), Inches(5.0), Inches(11.9), Inches(1.5),
                fill=LIGHT, line=GOLD, line_w=1.2, adj=0.07)
        _, tf = textbox(slide, Inches(0.95), Inches(5.22), Inches(11.2), Inches(1.1))
        write(tf, "The pattern: " + pat.get("en", ""), size=12.5, color=NAVY,
              bold=True, first=True, space_after=6)
        if pat.get("ar"):
            write(tf, pat["ar"], size=11.5, color=INK, arabic=True)
    footer_note(slide, "Ranked by engagement rate (engagement ÷ reach), not raw engagement.",
                "مرتبة حسب معدل التفاعل، لا حسب إجمالي التفاعل.")
    return slide


def slide_weak_posts(prs, d):
    slide = base_slide(prs)
    header(slide, "Underperforming Posts", "المنشورات الأقل أداءً", 7)
    ig = (d.get("ilaf") or {}).get("instagram") or {}
    posts = (ig.get("weak_posts") or [])[:3]
    if not posts:
        _, tf = textbox(slide, Inches(0.6), Inches(3.3), Inches(12.1), Inches(0.6))
        write(tf, NA, size=15, color=MUTED, align=PP_ALIGN.CENTER, first=True)
        return slide
    x, w = Inches(0.6), Inches(3.9)
    for i, post in enumerate(posts):
        _post_card(slide, x + i * (w + Inches(0.2)), Inches(1.5), w, Inches(3.25), post, NAVY)
    pat = ig.get("weak_pattern") or {}
    if pat.get("en") or pat.get("ar"):
        rounded(slide, Inches(0.6), Inches(5.0), Inches(11.9), Inches(1.5),
                fill=LIGHT, line=NAVY, line_w=1.2, adj=0.07)
        _, tf = textbox(slide, Inches(0.95), Inches(5.22), Inches(11.2), Inches(1.1))
        write(tf, "What to change: " + pat.get("en", ""), size=12.5, color=NAVY,
              bold=True, first=True, space_after=6)
        if pat.get("ar"):
            write(tf, pat["ar"], size=11.5, color=INK, arabic=True)
    return slide


def slide_google_business(prs, d):
    slide = base_slide(prs)
    header(slide, "Google Business Profile", "الملف التجاري على جوجل", 8)
    gb = (d.get("ilaf") or {}).get("google_business") or {}
    cur, prev = (gb.get("current") or {}), (gb.get("previous") or {})

    intent = [
        ("Phone calls", "المكالمات", cur.get("phone_calls"), prev.get("phone_calls")),
        ("Direction requests", "طلبات الاتجاهات", cur.get("direction_requests"),
         prev.get("direction_requests")),
        ("Website clicks", "النقرات على الموقع", cur.get("website_clicks"),
         prev.get("website_clicks")),
    ]
    x = Inches(0.6)
    for le, la, c, p_ in intent:
        dtxt, pct = delta(c, p_)
        kpi_tile(slide, x, Inches(1.55), Inches(3.9), Inches(1.95), le, la,
                 fmt_num(c) if has(c) else NA, dtxt, delta_color(pct, "up"))
        x += Inches(4.1)

    rows = [
        ("Total profile views", "إجمالي المشاهدات", "profile_views", "up"),
        ("Views on Search", "المشاهدات في البحث", "views_search", "up"),
        ("Views on Maps", "المشاهدات في الخرائط", "views_maps", "up"),
        ("Direct searches", "بحث مباشر", "searches_direct", "up"),
        ("Discovery searches", "بحث استكشافي", "searches_discovery", "up"),
    ]
    y = 3.72
    for le, la, key, good in rows:
        _metric_row(slide, y, le, la, cur.get(key), prev.get(key), fmt_num, good)
        y += 0.58

    rev = gb.get("reviews") or {}
    rounded(slide, Inches(7.0), Inches(3.75), Inches(5.5), Inches(2.75),
            fill=LIGHT, line=CREAM, line_w=1.2, adj=0.05)
    _, tf = textbox(slide, Inches(7.35), Inches(3.95), Inches(4.9), Inches(0.4))
    write(tf, "Reviews · التقييمات", size=12.5, color=NAVY, bold=True, first=True)
    _, tf2 = textbox(slide, Inches(7.35), Inches(4.4), Inches(4.9), Inches(0.5))
    rating = f"{rev['average_rating']} ★" if has(rev.get("average_rating")) else NA
    write(tf2, rating, size=22, color=MAGENTA, bold=True, first=True)
    _, tf3 = textbox(slide, Inches(7.35), Inches(4.95), Inches(4.9), Inches(0.35))
    write(tf3, f"{fmt_num(rev.get('total'))} total   ·   {fmt_num(rev.get('new_this_month'))} new this month",
          size=10.5, color=MUTED, first=True)
    if rev.get("quote_en") or rev.get("quote_ar"):
        _, tf4 = textbox(slide, Inches(7.35), Inches(5.35), Inches(4.9), Inches(1.0))
        if rev.get("quote_en"):
            write(tf4, f"“{rev['quote_en']}”", size=10, color=INK, italic=True,
                  first=True, space_after=4)
        if rev.get("quote_ar"):
            write(tf4, f"«{rev['quote_ar']}»", size=10, color=INK, italic=True, arabic=True)

    # Google retired the direct/discovery/branded search split in the newer Business Profile
    # performance API. Explaining a metric the slide is showing as unavailable reads as an
    # error, so only run that note when the figure is actually there.
    if has(cur.get("searches_discovery")):
        footer_note(slide,
                    "Discovery searches are people who found ILAF while looking for insurance — "
                    "the clearest measure of new market reach.",
                    "البحث الاستكشافي يمثل من وجد إيلاف أثناء بحثه عن التأمين.")
    else:
        footer_note(slide,
                    "Google no longer reports the direct / discovery search split, so those "
                    "rows show as unavailable rather than zero.",
                    "لم تعد جوجل تفصل بين البحث المباشر والاستكشافي، لذا تظهر هذه البنود كغير متاحة لا كصفر.")
    return slide


def slide_enquiries(prs, d):
    """Real enquiries: WhatsApp messages beside Google's calls and directions.

    Every other slide measures attention — reach, views, likes. This one measures
    people who actually made contact, which is the only number on the deck that maps
    directly to revenue. It belongs beside the Google figures, not on its own, because
    the comparison is the insight: which channel is actually producing conversations.
    """
    slide = base_slide(prs)
    header(slide, "Enquiries & Direct Contact", "الاستفسارات والتواصل المباشر")
    subtitle(slide,
             "People who actually got in touch this month — not views, not likes.",
             "من تواصل فعلياً هذا الشهر — لا مشاهدات ولا إعجابات.")

    wa = (d.get("ilaf") or {}).get("whatsapp") or {}
    wc, wp = (wa.get("current") or {}), (wa.get("previous") or {})
    gb = (d.get("ilaf") or {}).get("google_business") or {}
    gc, gp = (gb.get("current") or {}), (gb.get("previous") or {})

    tiles = [
        ("WhatsApp messages", "رسائل واتساب", wc.get("messages_received"),
         wp.get("messages_received"), MAGENTA),
        ("Phone calls (Google)", "المكالمات عبر جوجل", gc.get("phone_calls"),
         gp.get("phone_calls"), NAVY),
        ("Direction requests", "طلبات الاتجاهات", gc.get("direction_requests"),
         gp.get("direction_requests"), NAVY),
        ("Website clicks", "النقرات على الموقع", gc.get("website_clicks"),
         gp.get("website_clicks"), NAVY),
    ]
    x = Inches(0.6)
    for le, la, c, p_, _accent in tiles:
        dtxt, pct = delta(c, p_)
        kpi_tile(slide, x, Inches(1.95), Inches(2.98), Inches(2.0), le, la,
                 fmt_num(c) if has(c) else NA, dtxt, delta_color(pct, "up"))
        x += Inches(3.12)

    # Total direct contacts, only when every part of the sum is present. A partial
    # total reads as a real one and quietly understates the month.
    parts = [wc.get("messages_received"), gc.get("phone_calls")]
    prev_parts = [wp.get("messages_received"), gp.get("phone_calls")]
    rounded(slide, Inches(0.6), Inches(4.2), Inches(6.1), Inches(1.75),
            fill=LIGHT, line=GOLD, line_w=1.4, adj=0.07)
    _, tf = textbox(slide, Inches(0.95), Inches(4.4), Inches(5.4), Inches(0.36))
    write(tf, "Total conversations started", size=12, color=NAVY, bold=True, first=True)
    if all(has(v) for v in parts):
        total = sum(float(v) for v in parts)
        ptotal = sum(float(v) for v in prev_parts) if all(has(v) for v in prev_parts) else None
        dtxt, pct = delta(total, ptotal)
        _, tfv = textbox(slide, Inches(0.95), Inches(4.8), Inches(3.0), Inches(0.62))
        write(tfv, fmt_num(total), size=30, color=MAGENTA, bold=True, first=True)
        _, tfd = textbox(slide, Inches(4.0), Inches(5.0), Inches(2.4), Inches(0.4))
        write(tfd, dtxt or "—", size=12, color=delta_color(pct, "up"), bold=True, first=True)
    else:
        _, tfv = textbox(slide, Inches(0.95), Inches(4.85), Inches(5.2), Inches(0.5))
        write(tfv, NA, size=15, color=MUTED, bold=True, first=True)
    _, tfa = textbox(slide, Inches(0.95), Inches(5.5), Inches(5.4), Inches(0.34))
    write(tfa, "إجمالي المحادثات التي بدأت · WhatsApp + Google calls",
          size=10, color=MUTED, arabic=True, first=True)

    topics = wc.get("by_topic") or {}
    if topics:
        pairs = sorted(topics.items(), key=lambda kv: -float(kv[1] or 0))
        path = chart_hbar([k.replace("_", " ").title() for k, _ in pairs],
                          [float(v or 0) for _, v in pairs], "wa_topics.png",
                          highlight_index=0, xlabel="WhatsApp enquiries by topic")
        # Size by height, not width. The chart gains a row per topic, so a width-locked
        # picture gets taller as topics are added and eventually sits on top of the note.
        pic = slide.shapes.add_picture(path, Inches(7.0), Inches(4.15), height=Inches(1.8))
        if pic.width > Inches(5.5):
            pic.height = int(pic.height * Inches(5.5) / pic.width)
            pic.width = Inches(5.5)
    else:
        rounded(slide, Inches(7.0), Inches(4.2), Inches(5.5), Inches(1.75),
                fill=WHITE, line=CREAM, line_w=1.2, adj=0.07)
        _, tf2 = textbox(slide, Inches(7.35), Inches(4.45), Inches(4.9), Inches(1.3))
        write(tf2, "Breakdown by topic", size=12, color=NAVY, bold=True,
              first=True, space_after=6)
        write(tf2, "Not collected. Splitting WhatsApp enquiries by product — motor, "
                   "medical, travel, home — would show which product people actually "
                   "ask about, rather than which one we posted about.",
              size=9.5, color=MUTED)

    note = wa.get("note") or {}
    if note.get("en") or note.get("ar"):
        _, tf3 = textbox(slide, Inches(0.6), Inches(6.15), Inches(11.9), Inches(0.9))
        if note.get("en"):
            write(tf3, note["en"], size=11, color=NAVY, first=True, space_after=5)
        if note.get("ar"):
            write(tf3, note["ar"], size=10.5, color=INK, arabic=True)
    return slide


def slide_competitor_ads(prs, d):
    """Who is paying to be seen, and what they're pushing.

    Organic posts show what a competitor wants to say. Ads show where they are willing
    to spend, which is a far stronger signal of what they think is working.
    """
    slide = base_slide(prs)
    header(slide, "Competitor Advertising", "إعلانات المنافسين")
    subtitle(slide,
             "Sponsored ads competitors ran this month, from Meta's public Ad Library.",
             "الإعلانات المموّلة التي نشرها المنافسون هذا الشهر، من مكتبة إعلانات ميتا العامة.")

    ads = d.get("competitor_ads") or {}
    advertisers = ads.get("advertisers") or []

    if not advertisers:
        rounded(slide, Inches(1.6), Inches(2.4), Inches(10.1), Inches(2.5),
                fill=LIGHT, line=GOLD, line_w=1.5, adj=0.06)
        _, tf = textbox(slide, Inches(2.1), Inches(2.75), Inches(9.1), Inches(2.0))
        write(tf, NA, size=17, color=MUTED, bold=True, align=PP_ALIGN.CENTER,
              first=True, space_after=14)
        note = ads.get("note") or {}
        if note.get("en"):
            write(tf, note["en"], size=12, color=NAVY, align=PP_ALIGN.CENTER, space_after=10)
        if note.get("ar"):
            write(tf, note["ar"], size=11, color=INK, arabic=True, align=PP_ALIGN.CENTER)
        footer_note(slide,
                    "Meta's Ad Library only shows ads running at the moment you look, so "
                    "past months cannot be reconstructed — this is captured going forward.",
                    "تعرض مكتبة إعلانات ميتا الإعلانات النشطة وقت الاطلاع فقط، لذا يتعذّر "
                    "استرجاع الأشهر السابقة — ويبدأ الرصد من الآن.", y=6.9)
        return slide

    hdr = ["Advertiser", "Ads seen", "First seen", "What they promoted"]
    widths = [3.0, 1.2, 1.5, 6.2]
    x0, y0 = 0.6, 2.0
    rect(slide, Inches(x0), Inches(y0), Inches(sum(widths)), Inches(0.42), fill=NAVY)
    cx = x0
    for i, htxt in enumerate(hdr):
        _, tf = textbox(slide, Inches(cx + 0.12), Inches(y0 + 0.1), Inches(widths[i] - 0.2),
                        Inches(0.3))
        write(tf, htxt, size=10, color=WHITE, bold=True, first=True)
        cx += widths[i]

    y = y0 + 0.42
    for i, a in enumerate(advertisers[:8]):
        h = 0.62
        if i % 2 == 0:
            rect(slide, Inches(x0), Inches(y), Inches(sum(widths)), Inches(h), fill=LIGHT)
        cx = x0
        cells = [
            a.get("name_en") or a.get("handle") or "—",
            fmt_num(a.get("ads_count")),
            a.get("first_seen") or "—",
            a.get("promoted_en") or "—",
        ]
        for j, val in enumerate(cells):
            _, tf = textbox(slide, Inches(cx + 0.12), Inches(y + 0.08), Inches(widths[j] - 0.2),
                            Inches(h - 0.12))
            write(tf, str(val), size=9.5, color=INK,
                  bold=(j == 0), first=True)
            cx += widths[j]
        if a.get("promoted_ar"):
            _, tfa = textbox(slide, Inches(x0 + sum(widths[:3]) + 0.12), Inches(y + 0.33),
                             Inches(widths[3] - 0.2), Inches(0.26))
            write(tfa, a["promoted_ar"], size=8.5, color=MUTED, arabic=True, first=True)
        y += h

    ins = ads.get("insight") or {}
    if ins.get("en") or ins.get("ar"):
        rounded(slide, Inches(0.6), Inches(y + 0.2), Inches(11.9), Inches(1.15),
                fill=LIGHT, line=GOLD, line_w=1.3, adj=0.08)
        _, tf = textbox(slide, Inches(0.95), Inches(y + 0.38), Inches(11.2), Inches(0.9))
        if ins.get("en"):
            write(tf, ins["en"], size=11.5, color=NAVY, bold=True, first=True, space_after=5)
        if ins.get("ar"):
            write(tf, ins["ar"], size=11, color=INK, arabic=True)

    cd, td = ads.get("capture_days"), ads.get("capture_total_days")
    cov = (f"Captured on {cd} of {td} days. " if has(cd) and has(td) else "")
    footer_note(slide,
                cov + "Ad Library shows only currently-running ads, so a campaign that "
                      "started and ended between checks will not appear.",
                None, y=6.95)
    return slide


def slide_competitor_table(prs, d):
    slide = base_slide(prs)
    header(slide, "Competitor Activity", "نشاط المنافسين", 9)

    comps = d.get("competitors") or []
    ilaf_row = d.get("ilaf_competitor_row")
    rows = ([ilaf_row] if ilaf_row else []) + \
        [c for c in comps if c.get("type") == "takaful"] + \
        [c for c in comps if c.get("type") != "takaful"]
    rows = [r for r in rows if r]

    if not rows:
        _, tf = textbox(slide, Inches(0.6), Inches(3.3), Inches(12.1), Inches(0.6))
        write(tf, NA, size=15, color=MUTED, align=PP_ALIGN.CENTER, first=True)
        return slide

    headers = ["Account", "Type", "Followers", "Posts", "Static", "Carousel",
               "Reel", "Stories", "Posts/wk", "Rhythm"]
    widths = [2.9, 1.15, 1.2, 0.85, 0.85, 1.05, 0.8, 0.9, 0.95, 1.35]
    x0, y0 = Inches(0.55), Inches(1.5)
    rh = Inches(0.44)

    x = x0
    for h_, w_ in zip(headers, widths):
        rect(slide, x, y0, Inches(w_), rh, fill=NAVY)
        _, tf = textbox(slide, x + Inches(0.06), y0 + Inches(0.10), Inches(w_ - 0.12), Inches(0.3))
        write(tf, h_, size=9, color=WHITE, bold=True, first=True,
              align=PP_ALIGN.LEFT if h_ in ("Account", "Type", "Rhythm") else PP_ALIGN.RIGHT)
        x += Inches(w_)

    for ri, r in enumerate(rows):
        y = y0 + rh + ri * rh
        is_ilaf = bool(r.get("is_ilaf"))
        bg = CREAM if is_ilaf else (LIGHT if ri % 2 == 0 else WHITE)
        unavailable = r.get("availability") == "unavailable"
        bf = r.get("by_format") or {}
        st = r.get("stories") or {}
        cells = [
            r.get("handle", ""),
            "Takaful" if r.get("type") == "takaful" else "Conventional",
            fmt_num(r.get("followers")),
            fmt_num(r.get("posts_total")),
            fmt_num(bf.get("static")),
            fmt_num(bf.get("carousel")),
            fmt_num(bf.get("reel")),
            fmt_num(st.get("observed")),
            fmt_num(r.get("posts_per_week")),
            (r.get("rhythm") or "").title(),
        ]
        if unavailable:
            cells = [cells[0], cells[1]] + ["—"] * 8
        x = x0
        for ci, (c_, w_) in enumerate(zip(cells, widths)):
            rect(slide, x, y, Inches(w_), rh, fill=bg,
                 line=RGBColor(0xE4, 0xE0, 0xE8), line_w=0.5)
            _, tf = textbox(slide, x + Inches(0.06), y + Inches(0.11), Inches(w_ - 0.12), Inches(0.3))
            txt = c_ if c_ != NA else "—"
            write(tf, txt, size=9,
                  color=NAVY if is_ilaf else (MUTED if txt == "—" else INK),
                  bold=is_ilaf or ci == 0,
                  align=PP_ALIGN.LEFT if ci in (0, 1, 9) else PP_ALIGN.RIGHT, first=True)
            x += Inches(w_)

    partial = [r.get("handle") for r in rows if r.get("availability") in ("partial", "unavailable")]
    note_en = "ILAF's row is highlighted. Takaful peers are listed first as the like-for-like comparison."
    if partial:
        note_en += "  Incomplete data: " + ", ".join(partial) + "."
    footer_note(slide, note_en,
                "صف إيلاف مميّز، وشركات التكافل مدرجة أولاً لكونها المقارنة الأقرب.", y=6.85)
    return slide


def slide_strategy(prs, d):
    slide = base_slide(prs)
    header(slide, "Competitor Content Strategy", "استراتيجية محتوى المنافسين", 10)

    comps = d.get("competitors") or []
    ilaf_row = d.get("ilaf_competitor_row")
    entities, matrix = [], []
    all_themes = []
    pool = ([ilaf_row] if ilaf_row else []) + comps
    for r in pool:
        for t in (r.get("themes") or {}):
            if t not in all_themes:
                all_themes.append(t)
    for r in pool:
        th = r.get("themes") or {}
        if not th:
            continue
        entities.append((r.get("handle") or "").replace("@", ""))
        matrix.append([th.get(t, 0) for t in all_themes])

    if entities and all_themes:
        path = chart_stacked_themes(entities, all_themes, matrix, "themes.png")
        slide.shapes.add_picture(path, Inches(0.5), Inches(1.5), width=Inches(7.5))
    else:
        _, tf = textbox(slide, Inches(0.6), Inches(3.3), Inches(7.5), Inches(0.6))
        write(tf, NA, size=14, color=MUTED, align=PP_ALIGN.CENTER, first=True)

    obs = d.get("strategy_observations") or []
    _, tf = textbox(slide, Inches(8.3), Inches(1.55), Inches(4.5), Inches(5.0))
    write(tf, "What stands out", size=13, color=NAVY, bold=True, first=True, space_after=10)
    write(tf, "أبرز الملاحظات", size=11.5, color=MUTED, arabic=True, space_after=12)
    if not obs:
        write(tf, NA, size=11, color=MUTED)
    for o in obs[:5]:
        if o.get("en"):
            write(tf, "▸ " + o["en"], size=10.5, color=INK, space_after=3)
        if o.get("ar"):
            write(tf, o["ar"], size=10, color=MUTED, arabic=True, space_after=10)
    return slide


def slide_share_of_voice(prs, d):
    slide = base_slide(prs)
    header(slide, "Share of Voice", "الحصة من الصوت", 11)

    sov = d.get("share_of_voice") or {}
    share = sov.get("ilaf_share_pct")
    tracked = sov.get("accounts_tracked")
    total = sov.get("total_posts_market")
    ilaf_posts = ((d.get("ilaf_competitor_row") or {}).get("posts_total"))

    # "Share of voice" is a term of art. Anyone reading it cold — which is most of the
    # management audience — needs the arithmetic spelled out, or they will either
    # misread it or quote it wrongly elsewhere.
    if has(ilaf_posts) and has(total) and has(tracked):
        def_en = (f"ILAF published {fmt_num(ilaf_posts)} of the {fmt_num(total)} posts put out by "
                  f"the {tracked} tracked Kuwaiti insurers — this measures how much of the "
                  f"market's posting was ours.")
        def_ar = (f"نشرت إيلاف {fmt_num(ilaf_posts)} من أصل {fmt_num(total)} منشوراً لدى "
                  f"{tracked} شركات تأمين كويتية مرصودة — أي حصتنا من حجم النشر في السوق.")
    else:
        def_en = "How much of the Kuwaiti insurance market's Instagram posting came from ILAF."
        def_ar = "نسبة ما نشرته إيلاف من إجمالي نشاط شركات التأمين الكويتية على إنستغرام."
    subtitle(slide, def_en, def_ar)

    comps = d.get("competitors") or []
    ilaf_row = d.get("ilaf_competitor_row")
    pool = ([ilaf_row] if ilaf_row else []) + comps
    pairs = [((r.get("handle") or "").replace("@", ""), float(r.get("posts_total") or 0))
             for r in pool if has(r.get("posts_total"))]
    pairs.sort(key=lambda t: -t[1])

    note = sov.get("note") or {}

    # A share-of-voice chart needs a market to compare against. With only ILAF's own
    # posts counted, a donut or a single bar isn't an understated chart — it's a wrong
    # one, because it invites the reader to see a share where no denominator exists.
    # Say plainly that the measure isn't available yet instead.
    if not has(share) or len(pairs) < 2:
        rounded(slide, Inches(1.8), Inches(2.3), Inches(9.7), Inches(2.6),
                fill=LIGHT, line=GOLD, line_w=1.5, adj=0.06)
        _, tf = textbox(slide, Inches(2.3), Inches(2.65), Inches(8.7), Inches(2.0))
        write(tf, NA, size=17, color=MUTED, bold=True, align=PP_ALIGN.CENTER,
              first=True, space_after=14)
        if note.get("en"):
            write(tf, note["en"], size=12.5, color=NAVY, align=PP_ALIGN.CENTER, space_after=10)
        if note.get("ar"):
            write(tf, note["ar"], size=11.5, color=INK, arabic=True, align=PP_ALIGN.CENTER)
        if pairs:
            footer_note(slide,
                        f"ILAF published {int(pairs[0][1])} posts this month. A share figure "
                        f"needs the rest of the market's counts alongside it.", None, y=6.9)
        return slide

    path = chart_donut(share, "sov.png",
                       f"of {tracked} tracked accounts" if has(tracked) else "share of posts")
    slide.shapes.add_picture(path, Inches(1.35), Inches(2.15), height=Inches(3.5))

    ilaf_handle = ((ilaf_row or {}).get("handle") or "@ilaf_takaful").replace("@", "")
    hi = next((i for i, (h_, _) in enumerate(pairs) if h_ == ilaf_handle), -1)
    path = chart_hbar([p[0] for p in pairs], [p[1] for p in pairs], "sov_bars.png",
                      highlight_index=hi, xlabel="Posts published in the month")
    slide.shapes.add_picture(path, Inches(6.3), Inches(1.95), width=Inches(6.4))

    _, tf = textbox(slide, Inches(0.7), Inches(5.85), Inches(5.2), Inches(1.1))
    if note.get("en"):
        write(tf, note["en"], size=11.5, color=NAVY, bold=True, first=True, space_after=6)
    if note.get("ar"):
        write(tf, note["ar"], size=11, color=INK, arabic=True)

    if has(tracked):
        footer_note(slide, f"Share of all feed posts published by the {tracked} tracked "
                           f"Kuwaiti insurance accounts during the reporting window.",
                    None, y=6.95)
    return slide


def est_lines(text, width_in, size_pt):
    """Roughly how many lines this text will wrap to at the given width.

    Cards were previously a fixed height, which clipped any recommendation longer than
    two lines — the reader saw a sentence cut off mid-word and had no way to know what
    was missing. Estimating first and sizing the card to fit is the fix. Arabic sets
    tighter than Latin at the same point size, hence the two ratios.
    """
    if not text:
        return 0
    arabic = any("؀" <= c <= "ۿ" for c in text)
    char_w = size_pt * (0.46 if arabic else 0.50) / 72.0
    per_line = max(10, int(width_in / char_w))
    return max(1, -(-len(text) // per_line))


def _rec_metrics(r, width=11.9):
    """Line counts and total height for one recommendation card.

    The fitting loop and the drawing code both read from here. When each did its own
    arithmetic they drifted, and a card that "fitted" got drawn taller than the space
    reserved for it.
    """
    inner = width - 1.05
    l_en = est_lines(r.get("en"), inner, 11)
    l_ar = est_lines(r.get("ar"), inner, 10.5)
    l_ex = est_lines(r.get("example_en"), inner - 0.9, 10)
    l_exa = est_lines(r.get("example_ar"), inner - 0.9, 9.5)
    h = 0.16 + l_en * 0.185 + l_ar * 0.18
    if l_ex or l_exa:
        h += 0.2 + l_ex * 0.165 + l_exa * 0.16
    h += 0.1
    return h, l_en, l_ar, l_ex, l_exa, inner


def _rec_card(slide, y, n, r, width=11.9):
    """One recommendation, sized to its own content, with a worked example."""
    h, l_en, l_ar, l_ex, l_exa, inner = _rec_metrics(r, width)

    rounded(slide, Inches(0.6), Inches(y), Inches(width), Inches(h),
            fill=LIGHT, line=None, adj=0.10)
    rect(slide, Inches(0.6), Inches(y), Pt(5), Inches(h), fill=MAGENTA)
    _, tfn = textbox(slide, Inches(0.8), Inches(y + 0.16), Inches(0.45), Inches(0.4))
    write(tfn, str(n), size=15, color=MAGENTA, bold=True, first=True)

    ty = y + 0.11
    _, tfr = textbox(slide, Inches(1.35), Inches(ty), Inches(inner), Inches(l_en * 0.19 + 0.05))
    write(tfr, r.get("en", ""), size=11, color=INK, bold=True, first=True)
    ty += l_en * 0.185 + 0.02
    if r.get("ar"):
        _, tfa = textbox(slide, Inches(1.35), Inches(ty), Inches(inner), Inches(l_ar * 0.19 + 0.05))
        write(tfa, r["ar"], size=10.5, color=MUTED, arabic=True, first=True)
        ty += l_ar * 0.18

    # The example is the part people actually act on — a rule without a worked
    # instance gets nodded at and ignored.
    if r.get("example_en") or r.get("example_ar"):
        ty += 0.09
        rect(slide, Inches(1.35), Inches(ty), Inches(inner), Pt(0.75), fill=CREAM)
        ty += 0.07
        _, tfe = textbox(slide, Inches(1.35), Inches(ty), Inches(0.85), Inches(0.25))
        write(tfe, "Example", size=9, color=MAGENTA, bold=True, italic=True, first=True)
        if r.get("example_en"):
            _, tfx = textbox(slide, Inches(2.25), Inches(ty), Inches(inner - 0.9),
                             Inches(l_ex * 0.17 + 0.05))
            write(tfx, r["example_en"], size=10, color=NAVY, first=True)
        if r.get("example_ar"):
            _, tfxa = textbox(slide, Inches(2.25), Inches(ty + l_ex * 0.165),
                              Inches(inner - 0.9), Inches(l_exa * 0.17 + 0.05))
            write(tfxa, r["example_ar"], size=9.5, color=MUTED, arabic=True, first=True)
    return h


def slide_recommendations(prs, d):
    """Recommendations, paginated so nothing is ever clipped."""
    recs = d.get("recommendations") or []
    if not recs:
        slide = base_slide(prs)
        header(slide, "Recommendations", "التوصيات")
        _, tfe = textbox(slide, Inches(0.6), Inches(3.2), Inches(12.1), Inches(0.5))
        write(tfe, NA, size=15, color=MUTED, align=PP_ALIGN.CENTER, first=True)
        return slide

    BOTTOM = 6.95   # the footer rule sits at 7.0; anything above it is usable
    idx, page = 0, 0
    while idx < len(recs):
        page += 1
        slide = base_slide(prs)
        title = "Recommendations" + (f" ({page})" if page > 1 or len(recs) > 4 else "")
        header(slide, title, "التوصيات")
        _, tf = textbox(slide, Inches(0.6), Inches(1.28), Inches(11.9), Inches(0.34))
        write(tf, "What to do next month · ما يجب فعله الشهر القادم",
              size=12.5, color=NAVY, bold=True, first=True)
        y = 1.72
        while idx < len(recs):
            # measure first: a card that won't fit starts the next page instead of
            # spilling off the bottom of this one
            need = _rec_metrics(recs[idx])[0]
            if y + need > BOTTOM and y > 1.75:
                break
            y += _rec_card(slide, y, idx + 1, recs[idx]) + 0.13
            idx += 1
    return slide


def slide_coverage(prs, d):
    """What we could and could not measure — stated plainly, with an example."""
    slide = base_slide(prs)
    header(slide, "Data Coverage", "تغطية البيانات")
    subtitle(slide,
             "What this report could measure this month, and what it could not.",
             "ما تمكّن هذا التقرير من قياسه هذا الشهر، وما لم يتمكّن.")

    cov = d.get("coverage") or {}
    lines = cov.get("lines") or []

    # Legend. "Missing" is the one people misread, so it gets spelled out.
    legend = [("●", GOOD, "Measured", "Full month of data"),
              ("◐", GOLD, "Partial", "Some of the month, or some accounts"),
              ("○", BAD, "Not measured", "No data — this is NOT the same as zero")]
    lx = 0.6
    for mark, col, name, desc in legend:
        rounded(slide, Inches(lx), Inches(1.95), Inches(3.85), Inches(0.72),
                fill=LIGHT, line=None, adj=0.12)
        _, tf = textbox(slide, Inches(lx + 0.2), Inches(2.06), Inches(0.3), Inches(0.3))
        write(tf, mark, size=14, color=col, bold=True, first=True)
        _, tf2 = textbox(slide, Inches(lx + 0.58), Inches(2.05), Inches(3.1), Inches(0.55))
        write(tf2, name, size=11, color=NAVY, bold=True, first=True, space_after=1)
        write(tf2, desc, size=9, color=MUTED)
        lx += 4.02

    # The list grows each time a new measure joins the report, so the row pitch comes
    # from the count. A fixed pitch pushed the panel off the bottom of the slide the
    # moment the list reached eight entries.
    TOP, ROWS_END = 2.9, 5.72
    pitch = min(0.46, (ROWS_END - TOP) / max(1, len(lines)))
    tight = pitch < 0.40
    y = TOP
    for ln in lines:
        status = (ln.get("status") or "ok").lower()
        mark = {"ok": "●", "partial": "◐", "missing": "○"}.get(status, "●")
        col = {"ok": GOOD, "partial": GOLD, "missing": BAD}.get(status, MUTED)
        _, tf = textbox(slide, Inches(0.7), Inches(y), Inches(0.3), Inches(0.28))
        write(tf, mark, size=12 if tight else 13, color=col, bold=True, first=True)
        _, tf2 = textbox(slide, Inches(1.1), Inches(y + 0.01), Inches(4.6), Inches(0.28))
        write(tf2, ln.get("label", ""), size=10.5 if tight else 11, color=INK,
              bold=True, first=True)
        _, tf3 = textbox(slide, Inches(5.9), Inches(y + 0.02), Inches(6.6), Inches(0.28))
        write(tf3, ln.get("detail", ""), size=9.5 if tight else 10, color=MUTED, first=True)
        rect(slide, Inches(0.7), Inches(y + pitch - 0.11), Inches(11.8), Pt(0.5), fill=LIGHT)
        y += pitch

    ex = cov.get("example") or {}
    if ex.get("en") or ex.get("ar"):
        panel_y = min(y + 0.1, 5.8)
        rounded(slide, Inches(0.6), Inches(panel_y), Inches(11.9), Inches(1.0),
                fill=LIGHT, line=GOLD, line_w=1.3, adj=0.09)
        _, tf = textbox(slide, Inches(0.95), Inches(panel_y + 0.15), Inches(11.2), Inches(0.8))
        if ex.get("en"):
            write(tf, "Why this matters: " + ex["en"], size=11, color=NAVY,
                  bold=True, first=True, space_after=5)
        if ex.get("ar"):
            write(tf, ex["ar"], size=10.5, color=INK, arabic=True)

    note = cov.get("note") or {}
    if note.get("en"):
        footer_note(slide, note["en"], note.get("ar"), y=6.95)
    return slide


# --------------------------------------------------------------------------------------

BUILDERS = [
    slide_cover, slide_glance, slide_audience, slide_reach, slide_output,
    slide_top_posts, slide_weak_posts, slide_google_business, slide_enquiries,
    slide_competitor_table, slide_strategy, slide_competitor_ads, slide_share_of_voice,
    slide_recommendations, slide_coverage,
]


def build(data, out_path):
    global DATA_STATUS
    DATA_STATUS = (data.get("data_status") or "verified").lower()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    PAGE["n"] = 1
    PAGE["total"] = None            # first pass: discover the real slide count
    for fn in BUILDERS:
        try:
            fn(prs, data)
        except Exception as e:  # one bad slide shouldn't cost the whole deck
            sys.stderr.write(f"[warn] {fn.__name__} failed: {e}\n")
            slide = base_slide(prs)
            header(slide, "Slide unavailable", "الشريحة غير متاحة")
            _, tf = textbox(slide, Inches(0.6), Inches(3.2), Inches(12.1), Inches(0.8))
            write(tf, f"This slide could not be generated: {e}", size=12, color=MUTED,
                  align=PP_ALIGN.CENTER, first=True)
    # Recommendations paginate, so the total isn't known until everything is laid out.
    # Build once to count, then rebuild with the correct "n/total" on every slide.
    total = len(prs.slides._sldIdLst)
    if PAGE["total"] != total:
        PAGE["n"], PAGE["total"] = 1, total
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        for fn in BUILDERS:
            try:
                fn(prs, data)
            except Exception as e:
                sys.stderr.write(f"[warn] {fn.__name__} failed: {e}\n")
    prs.save(out_path)
    slim_pptx(out_path)
    return out_path



def slim_pptx(path):
    """Drop the template's unused slide layouts.

    python-pptx ships an 11-layout Office template; this deck only ever uses the blank
    one. The other ten travel in every file for no reason. Removing them is a large
    share of the deck's size once the charts are optimised.
    """
    import re, zipfile, os
    try:
        z = zipfile.ZipFile(path)
        names = z.namelist()
        used = set()
        for n in names:
            if n.startswith("ppt/slides/_rels/"):
                for m in re.findall(rb"slideLayouts/(slideLayout\d+\.xml)", z.read(n)):
                    used.add(m.decode())
        layouts = [n for n in names if re.match(r"ppt/slideLayouts/slideLayout\d+\.xml$", n)]
        drop = {n for n in layouts if os.path.basename(n) not in used}
        if not drop:
            z.close(); return path
        master = z.read("ppt/slideMasters/slideMaster1.xml").decode("utf8")
        mrels = z.read("ppt/slideMasters/_rels/slideMaster1.xml.rels").decode("utf8")
        drop_rids = set()
        for m in re.finditer(r'<Relationship Id="([^"]+)"[^>]*Target="\.\./slideLayouts/(slideLayout\d+\.xml)"', mrels):
            if "ppt/slideLayouts/" + m.group(2) in drop:
                drop_rids.add(m.group(1))
        for rid in drop_rids:
            mrels = re.sub(r'<Relationship Id="%s".*?/>' % re.escape(rid), "", mrels)
            master = re.sub(r'<p:sldLayoutId[^>]*r:id="%s"\s*/>' % re.escape(rid), "", master)
        ct = z.read("[Content_Types].xml").decode("utf8")
        for n in drop:
            ct = ct.replace('<Override PartName="/%s" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>' % n, "")
        skip = drop | {"ppt/slideLayouts/_rels/%s.rels" % os.path.basename(n) for n in drop}
        tmp = path + ".slim"
        out = zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED, compresslevel=9)
        for n in names:
            if n in skip:
                continue
            if n == "[Content_Types].xml":
                out.writestr(n, ct)
            elif n == "ppt/slideMasters/slideMaster1.xml":
                out.writestr(n, master)
            elif n == "ppt/slideMasters/_rels/slideMaster1.xml.rels":
                out.writestr(n, mrels)
            else:
                out.writestr(n, z.read(n))
        out.close(); z.close()
        Presentation(tmp)          # refuse to ship a file that won't reopen
        os.replace(tmp, path)
    except Exception as e:
        sys.stderr.write("[warn] slim_pptx skipped: %s\n" % e)
    return path


def main():
    global CHART_DIR
    ap = argparse.ArgumentParser(description="Build the ILAF monthly social report deck.")
    ap.add_argument("data", help="Path to report data JSON")
    ap.add_argument("-o", "--output", default=None, help="Output .pptx path")
    ap.add_argument("--chart-dir", default=None, help="Where to write chart PNGs")
    args = ap.parse_args()

    with open(args.data, encoding="utf-8") as f:
        data = json.load(f)

    out = args.output
    if not out:
        label = ((data.get("period") or {}).get("start") or "report")[:7]
        out = f"ILAF_Social_Report_{label}.pptx"

    CHART_DIR = args.chart_dir or os.path.join(
        os.path.dirname(os.path.abspath(out)) or ".", "_charts")
    os.makedirs(CHART_DIR, exist_ok=True)

    build(data, out)
    print(out)


if __name__ == "__main__":
    main()
